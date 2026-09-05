#!/usr/bin/env python3
"""
generate_pitch_deck.py
======================
Generates the official IRO (Intelligent Recovery Orchestrator) pitch presentation
deck for the Razorpay AI Buildathon 2026.

Theme: Razorpay Brand Palette (Deep Navy, Razorpay Blue, Crisp White, Cyan Accents)
Format: 16:9 Widescreen (13.333 x 7.5 inches)
Notes: Full pitch script embedded in PowerPoint presenter speaker notes.
"""

import os
import sys

# Ensure UTF-8 stdout encoding on Windows consoles
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# COLOR PALETTE (Razorpay Corporate & Fintech Identity)
# ==============================================================================
BG_DARK        = RGBColor(0x07, 0x0D, 0x1E)  # Deepest Navy background
CARD_BG        = RGBColor(0x0E, 0x17, 0x2E)  # Card container dark
CARD_BG_ALT    = RGBColor(0x13, 0x22, 0x47)  # Slightly elevated card
CARD_BG_DARK   = RGBColor(0x0B, 0x11, 0x24)  # Deep card inset
TEXT_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)  # Crisp pure white
TEXT_DIM       = RGBColor(0x94, 0xA3, 0xB8)  # Slate dim
TEXT_MUTED     = RGBColor(0x64, 0x74, 0x8B)  # Slate muted
BLUE_BRAND     = RGBColor(0x02, 0x84, 0xC7)  # Razorpay Signature Brand Blue
BLUE_VIBRANT   = RGBColor(0x25, 0x63, 0xEB)  # Electric Blue accent
BLUE_LIGHT     = RGBColor(0x60, 0xA5, 0xFA)  # Soft Blue for titles
CYAN_ACCENT    = RGBColor(0x06, 0xB6, 0xD4)  # High-tech Cyan/Teal
EMERALD_GREEN  = RGBColor(0x10, 0xB9, 0x81)  # Recovery Emerald
AMBER_WARN     = RGBColor(0xF5, 0x9E, 0x0B)  # Warning Amber
RED_DANGER     = RGBColor(0xEF, 0x44, 0x44)  # Payment Failure Red
BORDER_SUBTLE  = RGBColor(0x1E, 0x29, 0x3B)  # Card border subtle
BORDER_BLUE    = RGBColor(0x3B, 0x82, 0xF6)  # Active card border blue

FONT_MAIN = "Segoe UI"
FONT_CODE = "Consolas"

# ==============================================================================
# HELPER UTILITIES FOR SHAPES, TEXT & CARDS
# ==============================================================================
def create_blank_slide(prs):
    """Creates a blank slide with the default deep navy background."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    return slide

def add_header(slide, category, title, subtitle=None):
    """Adds a standardized top header banner with category badge, title and subtitle."""
    # Top Category Badge / Tracker
    badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
    tf = badge_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = category.upper()
    p.font.name = FONT_MAIN
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    # Main Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.55))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
    p_t = tf_t.paragraphs[0]
    p_t.text = title
    p_t.font.name = FONT_MAIN
    p_t.font.size = Pt(22)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE

    # Optional Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.35))
        tf_s = sub_box.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
        p_s = tf_s.paragraphs[0]
        p_s.text = subtitle
        p_s.font.name = FONT_MAIN
        p_s.font.size = Pt(12)
        p_s.font.color.rgb = TEXT_DIM

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER_SUBTLE):
    """Adds a stylized card container with solid fill and border."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1)
    return card

def add_speaker_notes(slide, notes_text):
    """Sets the presenter speaker notes text for the slide."""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

# ==============================================================================
# SLIDE 1: TITLE / HERO SLIDE
# ==============================================================================
def build_slide_1(prs):
    slide = create_blank_slide(prs)

    # Top Brand Pill
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(4.5), Inches(0.42))
    pill.fill.solid()
    pill.fill.fore_color.rgb = CARD_BG_ALT
    pill.line.color.rgb = BLUE_VIBRANT
    tf_pill = pill.text_frame
    tf_pill.margin_left = Inches(0.15)
    p_pill = tf_pill.paragraphs[0]
    p_pill.text = "RAZORPAY AI BUILDATHON 2026 • TRACK 3"
    p_pill.font.name = FONT_MAIN
    p_pill.font.size = Pt(10)
    p_pill.font.bold = True
    p_pill.font.color.rgb = BLUE_LIGHT

    # Main Project Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Intelligent Recovery Orchestrator"
    p.font.name = FONT_MAIN
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # Project Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(10.5), Inches(0.8))
    tf_tag = tagline_box.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "A distributed revenue-recovery system that detects failed payments, safely orchestrates recovery strategies, and knows when AI is worth using."
    p_tag.font.name = FONT_MAIN
    p_tag.font.size = Pt(15)
    p_tag.font.color.rgb = TEXT_DIM

    # 4 Quick Metric Badges
    stats = [
        ("57.0%", "RECOVERY RATE", "vs 32.0% Baseline", EMERALD_GREEN),
        ("₹231,069", "REVENUE SAVED", "+₹83,412 Net Lift", BLUE_LIGHT),
        ("0 UNSAFE", "ACTIONS EXECUTED", "24 Attempts Blocked", CYAN_ACCENT),
        ("2,920x", "AI COST ROI", "₹28.56 Total AI Cost", AMBER_WARN)
    ]

    card_w = Inches(2.7)
    card_h = Inches(1.5)
    start_x = Inches(0.8)
    gap_x = Inches(0.3)
    y_pos = Inches(4.1)

    for i, (val, label, subtext, color) in enumerate(stats):
        cx = start_x + i * (card_w + gap_x)
        add_card(slide, cx, y_pos, card_w, card_h, bg_color=CARD_BG, border_color=BORDER_SUBTLE)
        
        # Stat Value
        vbox = slide.shapes.add_textbox(cx + Inches(0.2), y_pos + Inches(0.18), card_w - Inches(0.4), Inches(0.5))
        tf_v = vbox.text_frame
        p_v = tf_v.paragraphs[0]
        p_v.text = val
        p_v.font.name = FONT_MAIN
        p_v.font.size = Pt(22)
        p_v.font.bold = True
        p_v.font.color.rgb = color

        # Label
        lbox = slide.shapes.add_textbox(cx + Inches(0.2), y_pos + Inches(0.72), card_w - Inches(0.4), Inches(0.3))
        tf_l = lbox.text_frame
        p_l = tf_l.paragraphs[0]
        p_l.text = label
        p_l.font.name = FONT_MAIN
        p_l.font.size = Pt(10)
        p_l.font.bold = True
        p_l.font.color.rgb = TEXT_WHITE

        # Subtext
        sbox = slide.shapes.add_textbox(cx + Inches(0.2), y_pos + Inches(1.02), card_w - Inches(0.4), Inches(0.3))
        tf_s = sbox.text_frame
        p_s = tf_s.paragraphs[0]
        p_s.text = subtext
        p_s.font.name = FONT_MAIN
        p_s.font.size = Pt(9)
        p_s.font.color.rgb = TEXT_MUTED

    # Bottom Metadata Bar
    meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.4))
    tf_m = meta_box.text_frame
    p_m = tf_m.paragraphs[0]
    p_m.text = "Built by Om Dapke  •  FastAPI + PostgreSQL + Redis Locks + Kafka Architecture + Multi-Tier AI Gateway"
    p_m.font.name = FONT_MAIN
    p_m.font.size = Pt(11)
    p_m.font.color.rgb = TEXT_MUTED

    add_speaker_notes(slide, 
        "A failed payment isn't necessarily a failed transaction. It's a decision problem.\n\n"
        "Do we retry? Switch the payment route? Ask the customer to take action? Wait for reconciliation? Or stop?\n\n"
        "And in payments, the wrong recovery decision can be worse than no recovery at all.\n\n"
        "So we built IRO — the Intelligent Recovery Orchestrator.\n\n"
        "Our first design decision was actually not to use AI everywhere."
    )

# ==============================================================================
# SLIDE 2: THE FUNDAMENTAL INSIGHT (5 RECOVERY DECISIONS)
# ==============================================================================
def build_slide_2(prs):
    slide = create_blank_slide(prs)
    add_header(slide, "01. The Problem Space", "The Fundamental Insight: 1 Failure → 5 Decisions",
               "In payments, the wrong recovery decision can be far worse than no recovery at all.")

    # Left: Central Payment Failed Simulation Card
    left_card = add_card(slide, Inches(0.8), Inches(1.9), Inches(3.6), Inches(4.5), bg_color=CARD_BG_ALT, border_color=RED_DANGER)
    
    # Text inside Left Card
    tbox = slide.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(3.0), Inches(3.9))
    tf = tbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "ORIGINATING EVENT"
    p.font.name = FONT_MAIN
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RED_DANGER

    p2 = tf.add_paragraph()
    p2.text = "₹75,000"
    p2.font.name = FONT_MAIN
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE

    p3 = tf.add_paragraph()
    p3.text = "PAYMENT FAILED\nAttempt 2 of 3 • Credit Card\nError: GATEWAY_TIMEOUT\nRoute: HDFC_CARDS (Degraded)"
    p3.font.name = FONT_CODE
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_DIM

    p4 = tf.add_paragraph()
    p4.text = "\nNaive Action:\nBlind retry on same route\nOutcome: Double Debit or Repeat Failure!"
    p4.font.name = FONT_MAIN
    p4.font.size = Pt(11)
    p4.font.color.rgb = RED_DANGER

    # Right: 5 Decision Paths Grid
    decisions = [
        ("1. RETRY?", "Try again with backoff", "Transient gateway glitch or network socket drop.", AMBER_WARN),
        ("2. SWITCH ROUTE?", "Try alternate payment rail", "Primary bank gateway degraded; route to healthy rail.", CYAN_ACCENT),
        ("3. ASK CUSTOMER?", "Interactive recovery link", "Expired 3DS OTP, insufficient balance, or card expired.", BLUE_LIGHT),
        ("4. WAIT?", "Hold for reconciliation", "Gateway timeout: bank may still be capturing in-flight.", BLUE_BRAND),
        ("5. STOP?", "Halt recovery immediately", "Hard decline, fraud stop, or merchant max retry reached.", RED_DANGER),
    ]

    card_w = Inches(3.7)
    card_h = Inches(1.2)
    start_x = Inches(4.7)
    
    for i, (title, action, rationale, col) in enumerate(decisions[:3]):
        y = Inches(1.9) + i * Inches(1.4)
        add_card(slide, start_x, y, card_w, card_h, bg_color=CARD_BG, border_color=BORDER_SUBTLE)
        tb = slide.shapes.add_textbox(start_x + Inches(0.2), y + Inches(0.12), card_w - Inches(0.4), card_h)
        tf_c = tb.text_frame
        p_c1 = tf_c.paragraphs[0]
        p_c1.text = title
        p_c1.font.name = FONT_MAIN
        p_c1.font.size = Pt(13)
        p_c1.font.bold = True
        p_c1.font.color.rgb = col
        p_c2 = tf_c.add_paragraph()
        p_c2.text = f"{action}\n{rationale}"
        p_c2.font.name = FONT_MAIN
        p_c2.font.size = Pt(10)
        p_c2.font.color.rgb = TEXT_DIM

    start_x_col2 = Inches(8.7)
    for i, (title, action, rationale, col) in enumerate(decisions[3:]):
        y = Inches(1.9) + i * Inches(1.4)
        add_card(slide, start_x_col2, y, card_w, card_h, bg_color=CARD_BG, border_color=BORDER_SUBTLE)
        tb = slide.shapes.add_textbox(start_x_col2 + Inches(0.2), y + Inches(0.12), card_w - Inches(0.4), card_h)
        tf_c = tb.text_frame
        p_c1 = tf_c.paragraphs[0]
        p_c1.text = title
        p_c1.font.name = FONT_MAIN
        p_c1.font.size = Pt(13)
        p_c1.font.bold = True
        p_c1.font.color.rgb = col
        p_c2 = tf_c.add_paragraph()
        p_c2.text = f"{action}\n{rationale}"
        p_c2.font.name = FONT_MAIN
        p_c2.font.size = Pt(10)
        p_c2.font.color.rgb = TEXT_DIM

    # Bottom Warning Strip
    warn_card = add_card(slide, Inches(8.7), Inches(4.7), card_w, Inches(1.7), bg_color=CARD_BG_DARK, border_color=AMBER_WARN)
    tb_w = slide.shapes.add_textbox(Inches(8.9), Inches(4.8), card_w - Inches(0.4), Inches(1.5))
    tf_w = tb_w.text_frame
    p_w1 = tf_w.paragraphs[0]
    p_w1.text = "THE COST OF GUESSING"
    p_w1.font.name = FONT_MAIN
    p_w1.font.size = Pt(10)
    p_w1.font.bold = True
    p_w1.font.color.rgb = AMBER_WARN
    p_w2 = tf_w.add_paragraph()
    p_w2.text = "• Double debits lead to chargeback penalties\n• Inappropriate retries trigger bank fraud throttles\n• Blind automation causes severe customer churn"
    p_w2.font.name = FONT_MAIN
    p_w2.font.size = Pt(10)
    p_w2.font.color.rgb = TEXT_DIM

    add_speaker_notes(slide,
        "Do we retry? Switch the payment route? Ask the customer to take action? Wait for reconciliation? Or stop?\n\n"
        "And in payments, the wrong recovery decision can be worse than no recovery at all.\n\n"
        "Most payment failures are predictable. A transient failure doesn't need an AI model deciding whether to retry.\n\n"
        "So IRO uses a decision hierarchy. Deterministic rules handle what we already understand, and AI is introduced only when the recovery decision becomes ambiguous."
    )

# ==============================================================================
# SLIDE 3: THE DECISION HIERARCHY (AI PROPOSES, SYSTEM DECIDES)
# ==============================================================================
def build_slide_3(prs):
    slide = create_blank_slide(prs)
    add_header(slide, "02. Architecture & Hierarchy", "The Decision Hierarchy: AI Proposes, System Decides",
               "Deterministic rules handle what we understand; AI is engaged only when decisions become ambiguous.")

    # 3 Execution Tiers
    tiers = [
        ("TIER 1: DETERMINISTIC ENGINE", "< 1ms Latency", "68% of All Traffic",
         "• Static rule-based evaluation & idempotent cache lookup\n• Automatically handles standard transient failures & terminal stops\n• Zero AI token overhead, deterministic SLA safety",
         CYAN_ACCENT),
        ("TIER 2: MODEL ROUTER & GATEWAY", "15ms Latency", "Classification & Routing",
         "• Intent classification for ambiguous errors\n• Selects optimal model tier (Fast Classifier vs Deep Specialist)\n• Evaluates route telemetry & merchant profile metadata",
         BLUE_LIGHT),
        ("TIER 3: SPECIALIST RECOVERY AGENT", "25ms Latency", "32% Deep Reasoning",
         "• Multi-step contextual investigation of payment history\n• Evaluates degraded rails, merchant SLA & customer risk\n• Proposes structured recovery plan (READ-ONLY permissions)",
         BLUE_VIBRANT),
    ]

    card_w = Inches(3.7)
    card_h = Inches(3.8)
    start_x = Inches(0.8)
    gap_x = Inches(0.3)
    y = Inches(1.9)

    for i, (title, lat, share, desc, col) in enumerate(tiers):
        cx = start_x + i * (card_w + gap_x)
        add_card(slide, cx, y, card_w, card_h, bg_color=CARD_BG, border_color=BORDER_SUBTLE)
        
        tb = slide.shapes.add_textbox(cx + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), card_h)
        tf = tb.text_frame
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_MAIN
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = f"{lat}  •  {share}"
        p2.font.name = FONT_CODE
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_DIM

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.name = FONT_MAIN
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = TEXT_WHITE

    # Bottom Card: Deterministic Financial Policy Guard (The Iron Gate)
    gate_card = add_card(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.1), bg_color=CARD_BG_ALT, border_color=EMERALD_GREEN)
    tb_g = slide.shapes.add_textbox(Inches(1.0), Inches(5.95), Inches(11.3), Inches(0.95))
    tf_g = tb_g.text_frame
    p_g1 = tf_g.paragraphs[0]
    p_g1.text = "MANDATORY ENFORCEMENT: DETERMINISTIC FINANCIAL POLICY ENGINE"
    p_g1.font.name = FONT_MAIN
    p_g1.font.size = Pt(11)
    p_g1.font.bold = True
    p_g1.font.color.rgb = EMERALD_GREEN

    p_g2 = tf_g.add_paragraph()
    p_g2.text = "Every AI recommendation must independently satisfy hard policy guards before execution: Retry limit thresholds, maximum recovery window SLA, merchant contract permissions, and atomic payment state revalidation. AI has NO direct write access."
    p_g2.font.name = FONT_MAIN
    p_g2.font.size = Pt(10)
    p_g2.font.color.rgb = TEXT_DIM

    add_speaker_notes(slide,
        "So IRO uses a decision hierarchy. Deterministic rules handle what we already understand, and AI is introduced only when the recovery decision becomes ambiguous.\n\n"
        "AI proposes. The system decides.\n\n"
        "Every AI recommendation passes through deterministic safety guards and a financial policy engine before any recovery action can execute."
    )

# ==============================================================================
# SLIDE 4: DISTRIBUTED CONCURRENCY & IDEMPOTENCY
# ==============================================================================
def build_slide_4(prs):
    slide = create_blank_slide(prs)
    add_header(slide, "03. Distributed Systems Engineering", "The Hard Part Wasn't AI: Concurrency & Idempotency",
               "In an event-driven payment system, duplicate events are inevitable. Duplicate debits are intolerable.")

    col_w = Inches(5.7)
    col_h = Inches(4.0)
    y = Inches(1.9)

    # Pillar 1: Redis Distributed Lock
    add_card(slide, Inches(0.8), y, col_w, col_h, bg_color=CARD_BG, border_color=RED_DANGER)
    tb1 = slide.shapes.add_textbox(Inches(1.1), y + Inches(0.25), col_w - Inches(0.6), col_h - Inches(0.5))
    tf1 = tb1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "1. REDIS DISTRIBUTED LOCKING"
    p1.font.name = FONT_MAIN
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = RED_DANGER

    p1_q = tf1.add_paragraph()
    p1_q.text = "Answers: \"Who can work on this payment now?\""
    p1_q.font.name = FONT_MAIN
    p1_q.font.size = Pt(12)
    p1_q.font.italic = True
    p1_q.font.color.rgb = TEXT_WHITE

    p1_desc = tf1.add_paragraph()
    p1_desc.text = (
        "\n• Prevents simultaneous recovery workers from retrying the same failure\n"
        "• Atomic key acquisition: SET lock:recovery:{payment_id} NX EX {ttl}\n"
        "• Automatic TTL expiration prevents deadlocks if a worker node crashes\n"
        "• Safe release via Lua script guarantees only lock owner can release\n"
        "• Critical for high-concurrency Kafka consumers processing failure events"
    )
    p1_desc.font.name = FONT_MAIN
    p1_desc.font.size = Pt(10.5)
    p1_desc.font.color.rgb = TEXT_DIM

    # Pillar 2: Database Durable Idempotency
    add_card(slide, Inches(6.8), y, col_w, col_h, bg_color=CARD_BG, border_color=CYAN_ACCENT)
    tb2 = slide.shapes.add_textbox(Inches(7.1), y + Inches(0.25), col_w - Inches(0.6), col_h - Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "2. DATABASE DURABLE IDEMPOTENCY"
    p2.font.name = FONT_MAIN
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = CYAN_ACCENT

    p2_q = tf2.add_paragraph()
    p2_q.text = "Answers: \"Has this recovery already happened?\""
    p2_q.font.name = FONT_MAIN
    p2_q.font.size = Pt(12)
    p2_q.font.italic = True
    p2_q.font.color.rgb = TEXT_WHITE

    p2_desc = tf2.add_paragraph()
    p2_desc.text = (
        "\n• Durable barrier against retried, redelivered, or delayed messages\n"
        "• Deterministic key hashing: idemp:rec:{payment_id}:{attempt_seq}\n"
        "• Pre-execution reservation pattern: record created before API call\n"
        "• If execution crashes mid-flight, state machine handles safe reconciliation\n"
        "• Guarantees EXACTLY-ONCE financial execution across network partitions"
    )
    p2_desc.font.name = FONT_MAIN
    p2_desc.font.size = Pt(10.5)
    p2_desc.font.color.rgb = TEXT_DIM

    # Bottom Synthesis Banner
    syn_card = add_card(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.9), bg_color=CARD_BG_ALT, border_color=BLUE_BRAND)
    tb_s = slide.shapes.add_textbox(Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.75))
    tf_s = tb_s.text_frame
    p_s = tf_s.paragraphs[0]
    p_s.text = "KEY FINTECH LESSON: A lock alone does not guarantee exactly-once payment execution. Distributed locking coordinates in-flight concurrency; relational idempotency provides durable historical finality."
    p_s.font.name = FONT_MAIN
    p_s.font.size = Pt(10.5)
    p_s.font.bold = True
    p_s.font.color.rgb = TEXT_WHITE

    add_speaker_notes(slide,
        "The hard part, however, wasn't just AI. It was everything around it.\n\n"
        "First was concurrency.\n\n"
        "In an event-driven payment system, the same failure can reach multiple workers. Two workers trying to recover the same payment could potentially create a duplicate financial action.\n\n"
        "So we use Redis for distributed locking, while the database provides the durable idempotency barrier.\n\n"
        "In simple terms, Redis answers who can work on this payment now, while the database answers whether this recovery has already happened."
    )

# ==============================================================================
# SLIDE 5: DYNAMIC STATE INVERSION & TIMEOUT AMBIGUITY
# ==============================================================================
def build_slide_5(prs):
    slide = create_blank_slide(prs)
    add_header(slide, "04. Financial Safety & Edge Cases", "Dynamic State Inversion & Timeout Ambiguity",
               "Payment recovery is an asynchronous time-travel problem where reality changes while you wait.")

    card_w = Inches(5.7)
    card_h = Inches(4.7)
    y = Inches(1.9)

    # Edge Case 1: Dynamic State Inversion
    add_card(slide, Inches(0.8), y, card_w, card_h, bg_color=CARD_BG, border_color=AMBER_WARN)
    tb1 = slide.shapes.add_textbox(Inches(1.1), y + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.6))
    tf1 = tb1.text_frame
    
    p1 = tf1.paragraphs[0]
    p1.text = "EDGE CASE 1: STATE INVERSION"
    p1.font.name = FONT_MAIN
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = AMBER_WARN

    p1_t = tf1.add_paragraph()
    p1_t.text = "Payment succeeds while recovery is waiting in queue."
    p1_t.font.name = FONT_MAIN
    p1_t.font.size = Pt(11)
    p1_t.font.italic = True
    p1_t.font.color.rgb = TEXT_WHITE

    p1_body = tf1.add_paragraph()
    p1_body.text = (
        "\n1. T0: Payment attempt fails on banking rail.\n"
        "2. T1: Recovery orchestrator schedules a retry for T + 60s.\n"
        "3. T2: Customer re-enters checkout and pays successfully via UPI!\n"
        "4. T3: Scheduled recovery worker wakes up to retry original payment.\n\n"
        "THE CATASTROPHIC RISK: Double debiting the customer!\n\n"
        "IRO DEFENSE MECHANISM: Authoritative Pre-Execution Revalidation.\n"
        "Before acquiring locks or initiating payment APIs, IRO queries the source-of-truth payment state. If status == SUCCESS, recovery is stopped immediately with zero action."
    )
    p1_body.font.name = FONT_MAIN
    p1_body.font.size = Pt(10)
    p1_body.font.color.rgb = TEXT_DIM

    # Edge Case 2: Ambiguous Timeouts
    add_card(slide, Inches(6.8), y, card_w, card_h, bg_color=CARD_BG, border_color=BLUE_VIBRANT)
    tb2 = slide.shapes.add_textbox(Inches(7.1), y + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.6))
    tf2 = tb2.text_frame

    p2 = tf2.paragraphs[0]
    p2.text = "EDGE CASE 2: TIMEOUT AMBIGUITY"
    p2.font.name = FONT_MAIN
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = BLUE_LIGHT

    p2_t = tf2.add_paragraph()
    p2_t.text = "A timeout does not mean the payment failed."
    p2_t.font.name = FONT_MAIN
    p2_t.font.size = Pt(11)
    p2_t.font.italic = True
    p2_t.font.color.rgb = TEXT_WHITE

    p2_body = tf2.add_paragraph()
    p2_body.text = (
        "\n1. Acquirer gateway drops connection after 15,000ms.\n"
        "2. Did the customer's bank debit the money? Unknown.\n"
        "3. Did the merchant capture the transaction? Unknown.\n\n"
        "THE CATASTROPHIC RISK: Blindly retrying creates a duplicate charge while the bank is still processing attempt #1.\n\n"
        "IRO DEFENSE MECHANISM: Hold for Reconciliation.\n"
        "IRO classifies timeouts as AMBIGUOUS_IN_FLIGHT. Instead of immediate retry, the case enters a holding state with automated polling against the bank settlement/status API."
    )
    p2_body.font.name = FONT_MAIN
    p2_body.font.size = Pt(10)
    p2_body.font.color.rgb = TEXT_DIM

    add_speaker_notes(slide,
        "The second challenge was that payment state can change while recovery is waiting.\n\n"
        "Imagine a payment fails, we schedule a recovery, and then the original payment succeeds.\n\n"
        "Before execution, IRO revalidates the current payment state. If the payment has already succeeded, the recovery is stopped.\n\n"
        "Another difficult case is a timeout. A timeout doesn't always mean the payment actually failed. The bank could still be processing it. Blindly retrying in that situation could create another financial action.\n\n"
        "So instead, IRO can hold the recovery and wait for reconciliation."
    )

# ==============================================================================
# SLIDE 6: AI AS INFRASTRUCTURE (TIERED ROUTING & L7 SCALING)
# ==============================================================================
def build_slide_6(prs):
    slide = create_blank_slide(prs)
    add_header(slide, "05. AI Infrastructure & Operations", "AI as Infrastructure: Tiered Routing & L7 Scaling",
               "The goal is not to use more AI. It is to use exactly as much intelligence as the decision requires.")

    layers = [
        ("LAYER 1: MODEL ROUTER (WHAT TO RUN?)",
         "Evaluates task complexity and routes to right-sized inference tiers.",
         "• Fast Classification Tier: Simple classification (timeout vs decline, route matching). 5ms latency, $0.0005/1k tokens.\n• Deep Reasoning Tier: Multi-attempt failure analysis, merchant contract reconciliation. 25ms latency, deep reasoning.",
         CYAN_ACCENT),
        ("LAYER 2: L7 RECOVERY LOAD BALANCER (WHERE TO RUN?)",
         "Health-aware dispatching across distributed specialist worker instances.",
         "• Active health monitoring: Tracks instance states (HEALTHY, DEGRADED, DOWN).\n• Telemetry-driven load balancing: Dispatches to least-loaded healthy specialist container.\n• Fast circuit breaker: Automatically trips if specialist latency exceeds 150ms SLA.",
         BLUE_LIGHT),
        ("LAYER 3: ZERO-AI DETERMINISTIC BYPASS (WHEN NOT TO RUN?)",
         "The most scalable AI request is the one you never make.",
         "• 68% of payment failures resolved with ZERO AI inference tokens.\n• Direct fallback: If all AI instances degrade or fail, system operates cleanly on deterministic rules.\n• Eliminates external API dependencies from the critical payment authorization path.",
         EMERALD_GREEN),
    ]

    card_h = Inches(1.55)
    gap_y = Inches(0.2)
    start_y = Inches(1.9)

    for i, (title, sub, details, col) in enumerate(layers):
        cy = start_y + i * (card_h + gap_y)
        add_card(slide, Inches(0.8), cy, Inches(11.7), card_h, bg_color=CARD_BG, border_color=BORDER_SUBTLE)
        
        tb = slide.shapes.add_textbox(Inches(1.1), cy + Inches(0.15), Inches(11.1), card_h - Inches(0.2))
        tf = tb.text_frame
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_MAIN
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = col

        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.name = FONT_MAIN
        p_sub.font.size = Pt(10)
        p_sub.font.italic = True
        p_sub.font.color.rgb = TEXT_WHITE

        p_det = tf.add_paragraph()
        p_det.text = details
        p_det.font.name = FONT_MAIN
        p_det.font.size = Pt(9.5)
        p_det.font.color.rgb = TEXT_DIM

    add_speaker_notes(slide,
        "And finally, once AI becomes part of the architecture, AI itself becomes an infrastructure problem.\n\n"
        "We use tiered model routing and an L7 load balancer so that simpler cases can use lightweight inference, while complex cases are routed to deeper reasoning capacity.\n\n"
        "The goal isn't to use more AI. It's to use exactly as much intelligence as the decision requires."
    )

# ==============================================================================
# SLIDE 7: CASE STUDY (THE ₹75,000 HIGH-VALUE RECOVERY)
# ==============================================================================
def build_slide_7(prs):
    slide = create_blank_slide(prs)
    add_header(slide, "06. End-to-End Walkthrough", "Case Study: The ₹75,000 High-Value Recovery",
               "Tracing an ambiguous edge case from failure detection to safe financial settlement.")

    steps = [
        ("1. FAILURE EVENT", "₹75,000 Payment Fails", "Credit card attempt #2 fails on HDFC rail with GATEWAY_ERROR. Rail health drops to 38%."),
        ("2. SPECIALIST INVESTIGATION", "Read-Only AI Analysis", "Ambiguity detected. Specialist agent inspects route health & merchant limits. (ZERO write access)."),
        ("3. AI PROPOSAL", "Recommend Route Switch", "Agent identifies ICICI rail is healthy (98.4%). Recommends SWITCH_ROUTE with 120s backoff."),
        ("4. POLICY ENGINE GATE", "Deterministic Validation", "Financial policy checks: Attempt limit (2<3), SLA window valid, route permitted, payment still failed."),
        ("5. SAFE EXECUTION", "Revenue Recovered!", "Redis TTL lock acquired, DB idempotency reserved, transaction captured on ICICI. ₹75,000 SAVED!"),
    ]

    card_w = Inches(2.18)
    card_h = Inches(4.6)
    gap_x = Inches(0.2)
    start_x = Inches(0.8)
    y = Inches(1.9)

    for i, (num, title, text) in enumerate(steps):
        cx = start_x + i * (card_w + gap_x)
        border_c = EMERALD_GREEN if i == 4 else (BLUE_VIBRANT if i == 2 else BORDER_SUBTLE)
        add_card(slide, cx, y, card_w, card_h, bg_color=CARD_BG, border_color=border_c)

        tb = slide.shapes.add_textbox(cx + Inches(0.15), y + Inches(0.2), card_w - Inches(0.3), card_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = num
        p1.font.name = FONT_MAIN
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = CYAN_ACCENT

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE

        p3 = tf.add_paragraph()
        p3.text = f"\n{text}"
        p3.font.name = FONT_MAIN
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = TEXT_DIM

    add_speaker_notes(slide,
        "Now let's see how this works in an actual recovery case.\n\n"
        "Imagine a seventy-five-thousand-rupee credit card payment that failed because its payment route has degraded, after multiple attempts.\n\n"
        "This isn't a straightforward retry anymore. IRO detects the ambiguity and sends the case to its specialist investigation layer.\n\n"
        "The agent can inspect the payment history, previous attempts, route health, and merchant recovery policy, but it has no write access to the payment system.\n\n"
        "Based on that evidence, the AI recommends switching to a healthier payment route. But that recommendation is not execution.\n\n"
        "The financial policy engine independently checks the retry limits, recovery window, merchant policy, permitted strategy, and current payment state.\n\n"
        "Only after those checks pass can execution happen. The concurrency lock is acquired, idempotency is reserved, and the recovery is executed.\n\n"
        "That's the difference between an AI agent acting on a payment and a payment system safely using AI."
    )

# ==============================================================================
# SLIDE 8: EMPIRICAL BENCHMARK EVALUATION
# ==============================================================================
def build_slide_8(prs):
    slide = create_blank_slide(prs)
    add_header(slide, "07. Empirical Evaluation", "Empirical Evaluation: 100 Synthetic Scenarios",
               "Head-to-head comparison of IRO vs. Naive Single-Rail Baseline on controlled synthetic payment data.")

    metrics = [
        ("RECOVERY RATE", "57.0%", "32.0%", "+78.1% Relative Lift", EMERALD_GREEN),
        ("RECOVERED REVENUE", "₹231,069", "₹147,657", "+₹83,412 Net Lift (+56.5%)", BLUE_LIGHT),
        ("UNSAFE ACTIONS", "0 EXECUTED", "24 EXECUTED", "24 Unsafe Attempts Blocked", CYAN_ACCENT),
        ("AI INFERENCE ROI", "2,920x", "0x (No AI)", "₹28.56 ($0.336) Total AI Cost", AMBER_WARN),
    ]

    card_w = Inches(2.7)
    card_h = Inches(2.1)
    gap_x = Inches(0.3)
    start_x = Inches(0.8)
    y = Inches(1.9)

    for i, (label, iro_val, base_val, lift, col) in enumerate(metrics):
        cx = start_x + i * (card_w + gap_x)
        add_card(slide, cx, y, card_w, card_h, bg_color=CARD_BG, border_color=BORDER_SUBTLE)

        tb = slide.shapes.add_textbox(cx + Inches(0.18), y + Inches(0.15), card_w - Inches(0.36), card_h)
        tf = tb.text_frame

        p_lbl = tf.paragraphs[0]
        p_lbl.text = label
        p_lbl.font.name = FONT_MAIN
        p_lbl.font.size = Pt(10)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = TEXT_MUTED

        p_iro = tf.add_paragraph()
        p_iro.text = f"IRO: {iro_val}"
        p_iro.font.name = FONT_MAIN
        p_iro.font.size = Pt(18)
        p_iro.font.bold = True
        p_iro.font.color.rgb = col

        p_base = tf.add_paragraph()
        p_base.text = f"Baseline: {base_val}"
        p_base.font.name = FONT_MAIN
        p_base.font.size = Pt(11)
        p_base.font.color.rgb = TEXT_DIM

        p_lift = tf.add_paragraph()
        p_lift.text = lift
        p_lift.font.name = FONT_MAIN
        p_lift.font.size = Pt(9.5)
        p_lift.font.bold = True
        p_lift.font.color.rgb = TEXT_WHITE

    table_card = add_card(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.8), bg_color=CARD_BG_ALT, border_color=BORDER_SUBTLE)
    tb_t = slide.shapes.add_textbox(Inches(1.1), Inches(4.45), Inches(11.1), Inches(1.5))
    tf_t = tb_t.text_frame

    p_th = tf_t.paragraphs[0]
    p_th.text = "BREAKDOWN BY TIER & EXECUTION SAFETY (N=100 TRANSACTIONS)"
    p_th.font.name = FONT_MAIN
    p_th.font.size = Pt(11)
    p_th.font.bold = True
    p_th.font.color.rgb = CYAN_ACCENT

    p_td = tf_t.add_paragraph()
    p_td.text = (
        "• Tier Distribution: 68 Tier-1 Deterministic (0 token overhead) + 32 Tier-3 Specialist Agent dispatches.\n"
        "• Token Consumption: 20,800 total tokens (14.4k prompt / 6.4k completion) across 32 complex dispatches.\n"
        "• Safety Invariant: Zero double debits, zero policy violations, zero unauthorized attempts executed."
    )
    p_td.font.name = FONT_MAIN
    p_td.font.size = Pt(10)
    p_td.font.color.rgb = TEXT_DIM

    disclaimer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.4))
    tf_d = disclaimer_box.text_frame
    p_d = tf_d.paragraphs[0]
    p_d.text = "DISCLAIMER: Evaluated across 100 controlled synthetic payment scenarios with realistic distributions. These are synthetic benchmark results, not real Razorpay transaction outcomes."
    p_d.font.name = FONT_MAIN
    p_d.font.size = Pt(9)
    p_d.font.italic = True
    p_d.font.color.rgb = TEXT_MUTED

    add_speaker_notes(slide,
        "We then evaluated IRO across one hundred controlled synthetic payment scenarios.\n\n"
        "IRO achieved a 57 percent overall recovery rate, compared with 32 percent for our transparent naive single-rail baseline.\n\n"
        "That resulted in two hundred and thirty-one thousand and sixty-nine rupees of simulated recovered revenue, compared with one hundred and forty-seven thousand six hundred and fifty-seven rupees for the baseline.\n\n"
        "And importantly, IRO executed zero unsafe actions while blocking twenty-four simulated unsafe recovery attempts.\n\n"
        "These are synthetic results, not real Razorpay transaction outcomes."
    )

# ==============================================================================
# SLIDE 9: CORE LEARNINGS & TAKEAWAYS
# ==============================================================================
def build_slide_9(prs):
    slide = create_blank_slide(prs)
    add_header(slide, "08. Key Takeaways", "Core Learnings: Payment Recovery Beyond AI",
               "Payment recovery is fundamentally a systems engineering challenge, and only sometimes an AI challenge.")

    pillars = [
        ("A RELIABILITY PROBLEM", "Resilient FSM state machines, dead letter queues, and bounded retry intervals.", CYAN_ACCENT),
        ("A CONCURRENCY PROBLEM", "Distributed Redis locking (SET NX EX) and durable relational idempotency barriers.", BLUE_LIGHT),
        ("A POLICY PROBLEM", "Deterministic financial policy guards, retry thresholds, and regulatory compliance boundaries.", AMBER_WARN),
        ("A COST PROBLEM", "Right-sized model routing, L7 load balancing, and fast deterministic bypass for 68% of cases.", RED_DANGER),
        ("AND ONLY SOMETIMES...", "An AI problem. Introduced strictly when recovery decisions become ambiguous.", EMERALD_GREEN),
    ]

    card_w = Inches(2.18)
    card_h = Inches(2.6)
    gap_x = Inches(0.2)
    start_x = Inches(0.8)
    y = Inches(1.9)

    for i, (title, desc, col) in enumerate(pillars):
        cx = start_x + i * (card_w + gap_x)
        add_card(slide, cx, y, card_w, card_h, bg_color=CARD_BG, border_color=BORDER_SUBTLE)

        tb = slide.shapes.add_textbox(cx + Inches(0.15), y + Inches(0.2), card_w - Inches(0.3), card_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_MAIN
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_DIM

    hero_card = add_card(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.0), bg_color=CARD_BG_ALT, border_color=BLUE_VIBRANT)
    tb_h = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(10.9), Inches(1.6))
    tf_h = tb_h.text_frame

    p_h1 = tf_h.paragraphs[0]
    p_h1.text = "\"We didn't build an AI that controls payments.\""
    p_h1.font.name = FONT_MAIN
    p_h1.font.size = Pt(16)
    p_h1.font.italic = True
    p_h1.font.color.rgb = TEXT_DIM

    p_h2 = tf_h.add_paragraph()
    p_h2.text = "We built a payment recovery system that knows when AI is worth using."
    p_h2.font.name = FONT_MAIN
    p_h2.font.size = Pt(24)
    p_h2.font.bold = True
    p_h2.font.color.rgb = TEXT_WHITE

    add_speaker_notes(slide,
        "Our biggest learning from building IRO was that payment recovery isn't just an AI problem.\n\n"
        "It's a reliability problem.\n"
        "A concurrency problem.\n"
        "A policy problem.\n"
        "A cost problem.\n"
        "And only sometimes… an AI problem.\n\n"
        "We didn't build an AI that controls payments.\n"
        "We built a payment recovery system that knows when AI is worth using.\n\n"
        "IRO — Intelligent Recovery Orchestrator."
    )

# ==============================================================================
# SLIDE 10: CONCLUSION & TECH STACK
# ==============================================================================
def build_slide_10(prs):
    slide = create_blank_slide(prs)
    add_header(slide, "09. Conclusion", "IRO: Production-Grade Revenue Recovery",
               "Full architectural verification, automated test suites, and open-source implementation.")

    add_card(slide, Inches(0.8), Inches(1.9), Inches(5.7), Inches(4.8), bg_color=CARD_BG, border_color=EMERALD_GREEN)
    tb_l = slide.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(4.3))
    tf_l = tb_l.text_frame

    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "VERIFICATION & PRODUCTION INVARIANTS"
    p_l1.font.name = FONT_MAIN
    p_l1.font.size = Pt(12)
    p_l1.font.bold = True
    p_l1.font.color.rgb = EMERALD_GREEN

    p_l2 = tf_l.add_paragraph()
    p_l2.text = (
        "\n• 120 / 120 Automated Tests Passing (Pytest)\n"
        "  Unit, Integration, FSM, Concurrency, and Gateway suites.\n\n"
        "• Strict Relational Audit Ledger\n"
        "  Every state transition and AI recommendation permanently logged.\n\n"
        "• Zero Code Mutations Needed for Scaling\n"
        "  Decoupled L7 load balancer handles multi-model worker pools.\n\n"
        "• Production-Grade Benchmark Framework\n"
        "  Empirical batch evaluation supporting 10 to 100,000 transactions."
    )
    p_l2.font.name = FONT_MAIN
    p_l2.font.size = Pt(10.5)
    p_l2.font.color.rgb = TEXT_DIM

    add_card(slide, Inches(6.8), Inches(1.9), Inches(5.7), Inches(4.8), bg_color=CARD_BG, border_color=BLUE_VIBRANT)
    tb_r = slide.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.3))
    tf_r = tb_r.text_frame

    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "CORE TECHNOLOGY STACK"
    p_r1.font.name = FONT_MAIN
    p_r1.font.size = Pt(12)
    p_r1.font.bold = True
    p_r1.font.color.rgb = CYAN_ACCENT

    p_r2 = tf_r.add_paragraph()
    p_r2.text = (
        "\n• Backend Framework: Python 3.12 • FastAPI • Pydantic v2\n"
        "• Persistence: PostgreSQL 16 (AsyncPG / SQLAlchemy 2.0)\n"
        "• Distributed State: Redis 7 (Atomic TTL Locks & Keys)\n"
        "• Event Backbone: Decoupled Kafka Event Contract Architecture\n"
        "• AI Tier: Tiered Model Router + L7 Adaptive Load Balancer\n"
        "• Observability: Real-Time Recovery Dashboard & Telemetry API\n\n"
        "PROJECT AUTHOR\n"
        "Om Dapke  •  github.com/omdapke01"
    )
    p_r2.font.name = FONT_MAIN
    p_r2.font.size = Pt(10.5)
    p_r2.font.color.rgb = TEXT_WHITE

    add_speaker_notes(slide,
        "IRO — Intelligent Recovery Orchestrator.\n\n"
        "The system is in charge of the AI, not the AI in charge of the payment system.\n\n"
        "Thank you."
    )

# ==============================================================================
# MAIN EXECUTION ENTRY POINT
# ==============================================================================
def main():
    print("Initializing PowerPoint presentation...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("Building Slide 1: Title / Hero...")
    build_slide_1(prs)

    print("Building Slide 2: The Fundamental Insight...")
    build_slide_2(prs)

    print("Building Slide 3: The Decision Hierarchy...")
    build_slide_3(prs)

    print("Building Slide 4: Concurrency & Idempotency...")
    build_slide_4(prs)

    print("Building Slide 5: Dynamic State & Timeout Ambiguity...")
    build_slide_5(prs)

    print("Building Slide 6: AI as Infrastructure (Tiered Routing & L7)...")
    build_slide_6(prs)

    print("Building Slide 7: Case Study (The INR 75,000 Recovery)...")
    build_slide_7(prs)

    print("Building Slide 8: Empirical Benchmark Evaluation...")
    build_slide_8(prs)

    print("Building Slide 9: Core Learnings & Takeaways...")
    build_slide_9(prs)

    print("Building Slide 10: Conclusion & Technology Stack...")
    build_slide_10(prs)

    output_path = "IRO_Razorpay_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation successfully saved to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    main()
